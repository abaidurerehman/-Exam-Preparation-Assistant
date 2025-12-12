# app.py - Exam Preparation Assistant (Hugging Face Deployment)
# ============================================================================

print("\n🚀 Starting Exam Preparation Assistant...")
print("=" * 70)

import os
import re
import json
import time
from typing import List, Tuple, Optional
from pathlib import Path
import warnings
import tempfile

import gradio as gr
import numpy as np

# Document processing
from PyPDF2 import PdfReader
from pptx import Presentation
import yt_dlp

# Embeddings and Vector Store
from sentence_transformers import SentenceTransformer
import faiss

# LangChain
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Gemini API
from google import genai
from google.genai import types

# Suppress warnings
warnings.filterwarnings("ignore")
# For env
from dotenv import load_dotenv
# ============================================================================
# CONFIGURATION & INITIALIZATION
# ============================================================================

# Get API key from environment variable (set in Hugging Face Spaces secrets)
load_dotenv()
API_KEY = os.getenv("GEMINI_API_KEY", "")
MODEL_NAME = "gemini-2.5-flash"
EMBEDDING_MODEL = "all-MiniLM-L6-v2"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
MAX_RETRIEVALS = 5

# Initialize Gemini Client
try:
    if API_KEY:
        client = genai.Client(api_key=API_KEY)
        print("✅ Gemini API initialized")
    else:
        print("⚠️ Warning: GEMINI_API_KEY not set.")
        print("📋 Please set GEMINI_API_KEY in Hugging Face Spaces secrets")
        client = None
except Exception as e:
    print(f"⚠️ Warning: {e}")
    client = None

# Initialize Embedding Model
print("Loading embedding model...")
embedding_model = SentenceTransformer(EMBEDDING_MODEL)
print("✅ Embedding model loaded\n")

# Global state
class DocumentState:
    def __init__(self):
        self.index = None
        self.documents = None
        self.embeddings = None
        self.metadata = None
        self.full_text = None

    def clear(self):
        self.index = None
        self.documents = None
        self.embeddings = None
        self.metadata = None
        self.full_text = None

    def is_loaded(self):
        return self.index is not None and self.documents is not None

doc_state = DocumentState()

# ============================================================================
# SAFETY & BIAS MITIGATION
# ============================================================================

SAFETY_FILTER_KEYWORDS = [
    "explicit", "violent", "hate", "illegal", "harmful",
    "inappropriate", "offensive", "adult", "nsfw"
]

def safety_check(text: str) -> Tuple[bool, str]:
    text_lower = text.lower()
    if len(text) < 5:
        return False, "Input too short. Please provide meaningful content."
    if len(text) > 50000:
        return False, "Input too long. Maximum 50,000 characters allowed."
    for keyword in SAFETY_FILTER_KEYWORDS:
        if keyword in text_lower:
            return False, f"Content contains potentially harmful keywords."
    if any(x in text.lower() for x in ["drop", "delete", "insert", "update", "exec"]):
        return False, "Input contains suspicious patterns."
    return True, "Safe"

# ============================================================================
# DOCUMENT PROCESSING
# ============================================================================

def extract_text_from_pdf(file_path: str) -> str:
    try:
        text = ""
        with open(file_path, 'rb') as pdf_file:
            pdf_reader = PdfReader(pdf_file)
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n--- Page {page_num + 1} ---\n" + page_text
        return text.strip()
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_pptx(file_path: str) -> str:
    try:
        text = ""
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_video_id(url: str) -> Optional[str]:
    """Extract video ID from YouTube URL"""
    patterns = [
        r'(?:youtube\.com\/watch\?v=|youtu\.be\/|youtube\.com\/embed\/)([^&\n?#]+)',
        r'youtube\.com\/watch\?v=([^&\n?#]+)',
        r'youtube\.com\/v\/([^&\n?#]+)',
        r'youtube\.com\/watch\?.*v=([^&\n?#]+)'
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None

def extract_text_from_youtube(url: str) -> str:
    """Extract content from YouTube - Enhanced version"""
    try:
        # Clean and validate URL
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        
        # Extract video ID
        video_id = extract_video_id(url)
        if not video_id:
            return f"❌ Could not extract video ID from URL: {url}"
        
        print(f"🔄 Processing YouTube video ID: {video_id}")
        
        # Configure yt-dlp options - minimal approach
        ydl_opts = {
            'quiet': True,
            'no_warnings': True,
            'extract_flat': False,
            'skip_download': True,
            'writesubtitles': False,
            'writeautomaticsub': False,
            'forcejson': True,
        }
        
        extracted_data = ""
        
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            try:
                print("📥 Fetching video information...")
                info = ydl.extract_info(url, download=False)
                
                # Extract comprehensive information
                title = info.get('title', 'Unknown Title')
                description = info.get('description', '')
                uploader = info.get('uploader', 'Unknown Channel')
                duration = info.get('duration', 0)
                view_count = info.get('view_count', 0)
                upload_date = info.get('upload_date', '')
                categories = info.get('categories', [])
                tags = info.get('tags', [])
                
                # Build rich content structure
                extracted_data += "=" * 60 + "\n"
                extracted_data += f"YOUTUBE VIDEO ANALYSIS\n"
                extracted_data += "=" * 60 + "\n\n"
                
                # Basic video info
                extracted_data += f"TITLE: {title}\n"
                extracted_data += f"CHANNEL: {uploader}\n"
                
                if duration:
                    hours = duration // 3600
                    minutes = (duration % 3600) // 60
                    seconds = duration % 60
                    if hours > 0:
                        extracted_data += f"DURATION: {hours}h {minutes}m {seconds}s\n"
                    else:
                        extracted_data += f"DURATION: {minutes}m {seconds}s\n"
                
                if view_count:
                    extracted_data += f"VIEWS: {view_count:,}\n"
                
                if upload_date:
                    # Format date: YYYYMMDD -> YYYY-MM-DD
                    formatted_date = f"{upload_date[:4]}-{upload_date[4:6]}-{upload_date[6:8]}"
                    extracted_data += f"UPLOAD DATE: {formatted_date}\n"
                
                extracted_data += "\n" + "-" * 40 + "\n\n"
                
                # Categories and tags
                if categories:
                    extracted_data += f"CATEGORIES: {', '.join(categories)}\n"
                
                if tags:
                    # Limit tags to avoid too much content
                    relevant_tags = tags[:10]  # First 10 tags
                    extracted_data += f"TAGS: {', '.join(relevant_tags)}\n"
                    if len(tags) > 10:
                        extracted_data += f"... and {len(tags) - 10} more tags\n"
                
                if categories or tags:
                    extracted_data += "\n" + "-" * 40 + "\n\n"
                
                # Description processing
                if description and len(description.strip()) > 50:
                    # Clean up description
                    clean_description = description.strip()
                    # Remove excessive whitespace
                    clean_description = re.sub(r'\n\s*\n', '\n\n', clean_description)
                    # Limit length
                    if len(clean_description) > 2000:
                        clean_description = clean_description[:2000] + "... [description truncated]"
                    
                    extracted_data += "VIDEO DESCRIPTION:\n"
                    extracted_data += clean_description + "\n"
                else:
                    extracted_data += "DESCRIPTION: No detailed description available.\n"
                    extracted_data += "You can still ask questions about this video's topic.\n"
                
                # Add video context based on title and metadata
                extracted_data += "\n" + "-" * 40 + "\n\n"
                extracted_data += "CONTEXT ANALYSIS:\n"
                extracted_data += f"Based on the title '{title}'"
                if categories:
                    extracted_data += f" in categories: {', '.join(categories[:3])}"
                if tags:
                    extracted_data += f" with relevant topics: {', '.join(tags[:5])}"
                extracted_data += "\n"
                
                print(f"✅ Successfully extracted metadata: {len(extracted_data)} characters")
                print(f"📹 Video: {title}")
                print(f"📝 Description length: {len(description)}")
                print(f"🏷️ Tags: {len(tags)}")
                print(f"📂 Categories: {categories}")
                
                return extracted_data
                
            except yt_dlp.DownloadError as e:
                error_msg = f"❌ YouTube download error: {str(e)}"
                print(error_msg)
                # Return basic info if available
                if 'title' in locals():
                    basic_content = f"TITLE: {title}\nCHANNEL: {uploader}\n\nNote: Could not extract full metadata. You can still ask questions about '{title}'."
                    return basic_content
                return f"❌ Error accessing video. The video might be private, age-restricted, or unavailable."
                
    except Exception as e:
        error_msg = f"❌ YouTube processing error: {str(e)}"
        print(error_msg)
        import traceback
        traceback.print_exc()
        return f"❌ Error processing YouTube URL. Please try a different video or check the URL."

# ============================================================================
# RAG PIPELINE
# ============================================================================

def chunk_text(text: str) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", " ", ""]
    )
    chunks = splitter.split_text(text)
    chunks = [c.strip() for c in chunks if c.strip()]
    return chunks

def create_embeddings(chunks: List[str]) -> bool:
    try:
        print(f"Creating embeddings for {len(chunks)} chunks...")
        embeddings = embedding_model.encode(chunks, batch_size=32, show_progress_bar=False)
        embeddings = np.array(embeddings).astype('float32')
        dimension = embeddings.shape[1]
        index = faiss.IndexFlatL2(dimension)
        index.add(embeddings)
        doc_state.index = index
        doc_state.documents = chunks
        doc_state.embeddings = embeddings
        print(f"✅ Created embeddings: {embeddings.shape}")
        return True
    except Exception as e:
        print(f"❌ Error creating embeddings: {e}")
        return False

def retrieve_relevant_chunks(query: str, k: int = MAX_RETRIEVALS) -> List[str]:
    try:
        if not doc_state.is_loaded():
            return []
        query_embedding = embedding_model.encode([query], show_progress_bar=False)
        query_embedding = np.array(query_embedding).astype('float32')
        distances, indices = doc_state.index.search(query_embedding, min(k, len(doc_state.documents)))
        retrieved = [doc_state.documents[idx] for idx in indices[0] if idx < len(doc_state.documents)]
        return retrieved
    except Exception as e:
        print(f"❌ Error retrieving chunks: {e}")
        return []

def process_document(file_obj) -> str:
    try:
        if file_obj is None:
            return "❌ No file selected"
        
        # Save uploaded file to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file_obj.name).suffix) as tmp_file:
            with open(file_obj.name, 'rb') as uploaded_file:
                tmp_file.write(uploaded_file.read())
            file_path = tmp_file.name

        print(f"Processing: {file_path}")
        if file_path.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith('.pptx'):
            text = extract_text_from_pptx(file_path)
        else:
            return "❌ Unsupported file type. Please upload PDF or PPTX."
        
        # Clean up temporary file
        try:
            os.unlink(file_path)
        except:
            pass
            
        if len(text) < 50 or "Error" in text:
            return "❌ Could not extract text from file."
        doc_state.full_text = text
        chunks = chunk_text(text)
        if len(chunks) == 0:
            return "❌ No content extracted."
        success = create_embeddings(chunks)
        if not success:
            return "❌ Error creating embeddings."
        doc_state.metadata = {
            "file": Path(file_obj.name).name,
            "chunks": len(chunks),
            "text_length": len(text),
        }
        msg = f"✅ Successfully processed: {Path(file_obj.name).name}\n📊 Chunks: {len(chunks)} | Size: {len(text)} chars"
        print(msg)
        return msg
    except Exception as e:
        error_msg = f"❌ Error: {str(e)}"
        print(error_msg)
        return error_msg

def process_youtube_url(url: str) -> str:
    try:
        if not url or len(url.strip()) < 10:
            return "❌ Please provide a valid YouTube URL"
        
        print(f"\n{'='*70}")
        print(f"🔄 Processing YouTube URL: {url}")
        print(f"{'='*70}")
        
        # Basic URL validation
        if 'youtube.com' not in url and 'youtu.be' not in url:
            return "❌ Please provide a valid YouTube URL (youtube.com or youtu.be)"
        
        text = extract_text_from_youtube(url)
        
        print(f"📊 Extracted text length: {len(text) if text else 0}")
        
        if not text or "❌" in text or len(text.strip()) < 100:
            print("❌ Extraction failed - insufficient content")
            return f"❌ Could not extract sufficient content from this video.\n\nPossible reasons:\n• Video is private or age-restricted\n• No description available\n• Network restrictions\n\nYou can try:\n• A different YouTube video\n• Uploading PDF/PPTX files instead\n• Using the Q&A feature with general questions"
        
        print("✅ Text extracted successfully, creating embeddings...")
        
        doc_state.full_text = text
        chunks = chunk_text(text)
        
        if len(chunks) == 0:
            print("❌ No chunks created")
            return "❌ No meaningful content extracted from the video."
        
        print(f"📦 Created {len(chunks)} chunks")
        
        success = create_embeddings(chunks)
        if not success:
            print("❌ Embedding creation failed")
            return "❌ Error creating embeddings."
        
        doc_state.metadata = {
            "source": "YouTube",
            "url": url,
            "chunks": len(chunks),
            "text_length": len(text),
        }
        msg = f"✅ Successfully processed YouTube video\n📊 Chunks: {len(chunks)} | Size: {len(text)} chars\n\nNote: Using video metadata, description, and context for analysis."
        print(msg)
        return msg
    except Exception as e:
        error_msg = f"❌ Error processing YouTube URL: {str(e)}"
        print(error_msg)
        import traceback
        traceback.print_exc()
        return f"❌ Error: {str(e)}\n\nPlease try a different video or upload PDF/PPTX files."

# ============================================================================
# CLEAN OUTPUT FORMATTING - NO BOLD MARKDOWN
# ============================================================================

def format_output(title: str, content: str, emoji: str = "📚") -> str:
    """Create clean formatted output without markdown"""
    separator = "=" * 60
    return f"\n{separator}\n{emoji} {title}\n{separator}\n\n{content}\n\n{separator}\n"

# ============================================================================
# GENERATION FUNCTIONS - CONCISE VERSION
# ============================================================================

CLEAN_STYLE_INSTRUCTIONS = """
Write in clear, concise paragraphs.
Use numbered sections with clear headings.
Do not use markdown formatting like **bold** or bullet points.
Provide structured, well-organized content.
Keep it academic but straightforward.
Be direct and to the point.
"""

def generate_summary() -> str:
    if not doc_state.is_loaded() or not client:
        return "❌ No document loaded or API not configured"
    try:
        context = doc_state.full_text if doc_state.full_text and len(doc_state.full_text) < 5000 else "\n\n".join(doc_state.documents[:10])
        prompt = f"""
{CLEAN_STYLE_INSTRUCTIONS}

Create a concise summary with these sections:

1. OVERVIEW
Brief introduction to the main topic

2. KEY CONCEPTS
List the 3-5 most important concepts discussed

3. MAIN POINTS
Explain each concept with examples

4. SUMMARY
Essential takeaways

Content to analyze:
{context}

Keep each section brief and focused.
"""
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.7,
                top_p=0.9,
                max_output_tokens=1500,
            )
        )
        result = response.text if response.text else "❌ No summary generated."
        return format_output("SUMMARY", result, "📝")
    except Exception as e:
        return f"❌ Error: {str(e)}"

def generate_practice_questions(num_questions: int = 5) -> str:
    if not doc_state.is_loaded() or not client:
        return "❌ No document loaded or API not configured"
    try:
        context = "\n\n".join(doc_state.documents[:min(15, len(doc_state.documents))])
        prompt = f"""
{CLEAN_STYLE_INSTRUCTIONS}

Generate exactly {num_questions} exam-style questions with clear answers.

Format each as:

QUESTION 1: [Question text here]
ANSWER 1: [Clear answer here, 2-3 sentences]

QUESTION 2: [Question text here]  
ANSWER 2: [Clear answer here, 2-3 sentences]

Requirements:
- Questions should test core concepts
- Answers should be direct and comprehensive
- Cover different topics from the material
- Keep questions and answers concise

Content:
{context}
"""
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.8,
                top_p=0.9,
                max_output_tokens=2000,
            )
        )
        result = response.text if response.text else "❌ No questions generated."
        return format_output("PRACTICE QUESTIONS & ANSWERS", result, "❓")
    except Exception as e:
        return f"❌ Error: {str(e)}"

def generate_study_plan() -> str:
    if not doc_state.is_loaded() or not client:
        return "❌ No document loaded or API not configured"
    try:
        context = "\n\n".join(doc_state.documents[:min(20, len(doc_state.documents))])
        prompt = f"""
{CLEAN_STYLE_INSTRUCTIONS}

Create a practical study plan with these sections:

1. TOPIC ORGANIZATION
How to structure the material

2. STUDY SCHEDULE  
Suggested timeline and daily focus

3. KEY AREAS
Most important concepts to master

4. STUDY METHODS
Effective learning techniques

5. PRACTICE APPROACH
How to test your knowledge

Content:
{context}

Provide actionable, straightforward advice.
"""
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.7,
                top_p=0.9,
                max_output_tokens=2000,
            )
        )
        result = response.text if response.text else "❌ No study plan generated."
        return format_output("STUDY PLAN", result, "📅")
    except Exception as e:
        return f"❌ Error: {str(e)}"

def answer_question(question: str) -> str:
    if not doc_state.is_loaded() or not client:
        return "❌ No document loaded or API not configured"
    is_safe, msg = safety_check(question)
    if not is_safe:
        return f"❌ {msg}"
    try:
        relevant_chunks = retrieve_relevant_chunks(question, k=MAX_RETRIEVALS)
        if not relevant_chunks:
            return "❌ No relevant information found in the document."
        context = "\n\n".join(relevant_chunks)
        prompt = f"""
{CLEAN_STYLE_INSTRUCTIONS}

Answer this question based on the provided content:

QUESTION: {question}

CONTENT:
{context}

Provide a clear, well-structured answer.
If the answer has multiple points, use numbered sections.
Be direct and avoid unnecessary elaboration.
"""
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.7,
                top_p=0.9,
                max_output_tokens=1200,
            )
        )
        result = response.text if response.text else "❌ No answer generated."
        return format_output(f"ANSWER: {question[:50]}...", result, "💡")
    except Exception as e:
        return f"❌ Error: {str(e)}"

# ============================================================================
# GRADIO INTERFACE - COMPATIBLE VERSION
# ============================================================================

print("Creating interface...")

with gr.Blocks() as app:
    
    # Custom CSS embedded in the interface
    gr.HTML("""
    <style>
    body, .gradio-container {
        background: #F4F5F7 !important;
        color: #1E1E1E !important;
        font-family: 'Segoe UI', Roboto, 'Helvetica Neue', sans-serif !important;
    }
    .gr-block { padding: 20px !important; }
    .gr-row { gap: 20px; }
    .card {
        background: #FFFFFF;
        border-radius: 12px;
        box-shadow: 0 6px 20px rgba(16, 24, 40, 0.06);
        padding: 20px;
        border: 1px solid rgba(15, 23, 42, 0.04);
    }
    .header { text-align: center; padding: 30px 12px; }
    h1, h2, h3 { color: #1E1E1E; font-weight: 700; margin: 0; }
    h1 { font-size: 28px; }
    .lead { margin-top: 10px; color: #5B6B77; font-size: 16px; }
    input[type="text"], textarea, .gr-file-upload {
        font-size: 14px !important;
        border-radius: 8px !important;
        padding: 12px !important;
        border: 1px solid rgba(16, 24, 40, 0.08) !important;
        background: #fff !important;
        color: #1E1E1E !important;
    }
    .gr-button, button {
        background: #000 !important;
        color: white !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 11px 18px !important;
        font-size: 15px !important;
        font-weight: 600 !important;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.12) !important;
        cursor: pointer !important;
    }
    .gr-button:hover, button:hover { background: #222 !important; }
    .gr-tabs .gr-tab { border-radius: 8px !important; padding: 12px !important; }
    .gr-textbox, textarea { font-family: 'Segoe UI', sans-serif !important; line-height: 1.7 !important; }
    </style>
    """)
    
    gr.HTML("""
    <div class="header card">
        <h1>📚 Exam Preparation Assistant</h1>
        <p class="lead">Upload study materials and get structured summaries, questions, and study plans</p>
    </div>
    """)

    with gr.Tabs():
        with gr.Tab("📤 Upload Documents"):
            with gr.Row():
                with gr.Column():
                    gr.Markdown("### 📄 Upload File")
                    file_input = gr.File(label="PDF or PPTX", file_types=[".pdf", ".pptx"])
                    upload_btn = gr.Button("Process Document", scale=1)
                    upload_status = gr.Textbox(label="Status", lines=4, interactive=False)
                
                with gr.Column():
                    gr.Markdown("### 🎥 YouTube URL")
                    gr.Markdown("""
**Supported URL Formats:**
- `https://www.youtube.com/watch?v=VIDEO_ID`
- `https://youtu.be/VIDEO_ID`
- `https://youtube.com/watch?v=VIDEO_ID`

**Note:** YouTube processing extracts metadata, description, and context. 
Some videos may have limited content available.

**Example URLs to test:**
- https://www.youtube.com/watch?v=dQw4w9WgXcQ
- https://youtu.be/dQw4w9WgXcQ
                    """)
                    youtube_url = gr.Textbox(label="Paste YouTube URL", placeholder="https://www.youtube.com/watch?v=...")
                    youtube_btn = gr.Button("Process YouTube", scale=1)
                    youtube_status = gr.Textbox(label="Status", lines=4, interactive=False)

        with gr.Tab("📝 Summary"):
            with gr.Row():
                summary_btn = gr.Button("🔄 Generate Summary", scale=1, size="lg")
            summary_output = gr.Textbox(label="Summary", lines=20, interactive=False)

        with gr.Tab("❓ Questions"):
            with gr.Row():
                questions_btn = gr.Button("🔄 Generate Questions", scale=1, size="lg")
            questions_output = gr.Textbox(label="Q&A", lines=20, interactive=False)

        with gr.Tab("📅 Study Plan"):
            with gr.Row():
                plan_btn = gr.Button("🔄 Generate Study Plan", scale=1, size="lg")
            plan_output = gr.Textbox(label="Study Plan", lines=20, interactive=False)

        with gr.Tab("💬 Ask Question"):
            user_question = gr.Textbox(label="Your Question", placeholder="Ask anything about the material...", lines=3)
            ask_btn = gr.Button("Get Answer", scale=1)
            answer_output = gr.Textbox(label="Answer", lines=15, interactive=False)

    # Event bindings
    upload_btn.click(process_document, inputs=file_input, outputs=upload_status)
    youtube_btn.click(process_youtube_url, inputs=youtube_url, outputs=youtube_status)
    summary_btn.click(generate_summary, outputs=summary_output)
    questions_btn.click(generate_practice_questions, outputs=questions_output)
    plan_btn.click(generate_study_plan, outputs=plan_output)
    ask_btn.click(answer_question, inputs=user_question, outputs=answer_output)

print("✅ Interface ready. Launching...\n")

if __name__ == "__main__":
    app.launch()